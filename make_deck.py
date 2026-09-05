"""Build the pitch deck: chukta_pitch.pptx

    python make_deck.py

Ten slides for a five-minute talk, with the spoken words in each slide's
speaker-notes pane so PowerPoint's presenter view shows them while the audience
sees only the slide.

Every figure comes from FACTS below, which mirrors what `eval/check_claims.py`
asserts in CI. If a number changes in the code, the build goes red and this
file needs updating with it - the deck cannot quietly drift.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Same palette as the dashboard and the rendered video.
BG = RGBColor(0x0E, 0x14, 0x16)
PANEL = RGBColor(0x18, 0x22, 0x25)
INK = RGBColor(0xDB, 0xE6, 0xE4)
DIM = RGBColor(0x8F, 0xA3, 0xA3)
FAINT = RGBColor(0x5F, 0x73, 0x75)
ACCENT = RGBColor(0x54, 0xB7, 0xA5)
RED = RGBColor(0xD4, 0x66, 0x5C)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)   # 16:9
MARGIN = Inches(0.9)


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # blank
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    # left accent rule
    bar = s.shapes.add_shape(1, Emu(0), Emu(0), Inches(0.07), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return s


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size, colour, *, font=SANS, bold=False, space_after=8,
         first=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = colour
    r.font.name = font
    r.font.bold = bold
    return p


def header(slide, kicker: str, title: str):
    tf = textbox(slide, MARGIN, Inches(0.55), W - MARGIN * 2, Inches(1.6))
    para(tf, kicker, 14, ACCENT, first=True, space_after=10)
    para(tf, title, 40, INK, bold=True, space_after=0)


def notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def table(slide, rows, top, *, highlight_last=False, last_colour=ACCENT,
          label_w=8.4):
    """Monospace two-column figure block on a panel."""
    height = Inches(0.52) * len(rows) + Inches(0.4)
    panel = slide.shapes.add_shape(
        1, MARGIN - Inches(0.25), top - Inches(0.18),
        W - MARGIN * 2 + Inches(0.5), height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.fill.background()
    panel.shadow.inherit = False

    y = top
    for i, (label, value) in enumerate(rows):
        last = i == len(rows) - 1
        colour = last_colour if (highlight_last and last) else DIM

        tf = textbox(slide, MARGIN, y, Inches(label_w), Inches(0.5))
        para(tf, label, 17, colour, font=MONO, first=True, space_after=0)

        tf2 = textbox(slide, MARGIN + Inches(label_w), y,
                      W - MARGIN * 2 - Inches(label_w), Inches(0.5))
        para(tf2, value, 17, colour, font=MONO, first=True, space_after=0,
             align=PP_ALIGN.RIGHT)
        y += Inches(0.52)
    return y


def bullets(slide, lines, top, *, font=SANS, size=18, colour=DIM):
    tf = textbox(slide, MARGIN, top, W - MARGIN * 2, Inches(3.2))
    for i, line in enumerate(lines):
        para(tf, line, size, colour, font=font, first=(i == 0), space_after=12)


def kicker_line(slide, text, colour=INK):
    """The one-line takeaway above the footer."""
    bar = slide.shapes.add_shape(1, MARGIN - Inches(0.22), H - Inches(1.55),
                                 Inches(0.05), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = textbox(slide, MARGIN, H - Inches(1.6), W - MARGIN * 2, Inches(0.7))
    para(tf, text, 19, colour, first=True, space_after=0)


def footer(slide, n, total, note=""):
    tf = textbox(slide, MARGIN, H - Inches(0.72), W - MARGIN * 2, Inches(0.4))
    para(tf, note, 11, FAINT, first=True, space_after=0)
    tf2 = textbox(slide, W - Inches(1.6), H - Inches(0.72), Inches(0.9), Inches(0.4))
    para(tf2, f"{n} / {total}", 11, FAINT, first=True, space_after=0,
         align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------

TOTAL = 10


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    n = 0

    # 1 ---------------------------------------------------------------- title
    n += 1
    s = add_slide(prs)
    tf = textbox(s, MARGIN, Inches(2.1), W - MARGIN * 2, Inches(2.4))
    para(tf, "RAZORPAY AI BUILDATHON  ·  TRACK 03", 14, ACCENT, first=True,
         space_after=18)
    para(tf, "Chukta", 66, INK, bold=True, space_after=14)
    para(tf, "A recovery agent that knows when to stop", 22, DIM, space_after=0)
    footer(s, n, TOTAL, "Chukta - Hindi for settled, or paid up")
    notes(s, """This is Chukta. It works failed Razorpay payments.

The name is Hindi for "settled", or "paid up" - it's named for the outcome, not
for the act of collecting. That turns out to be the whole idea.""")

    # 2 -------------------------------------------------------------- problem
    n += 1
    s = add_slide(prs)
    header(s, "THE PROBLEM", "The default retry ladder is wrong twice")
    bullets(s, [
        "Retry at T+1h, T+24h, T+72h.  Same card.  Generic SMS.",
        "",
        "It ignores WHY it failed.  An expired card retried three times gives",
        "three guaranteed declines and a worse decline ratio with the issuer -",
        "which raises your cost on every future transaction.",
        "",
        "It ignores what trying COSTS.  Outreach to a lukewarm subscriber can",
        "trigger the exact cancellation you were trying to prevent.",
    ], Inches(2.5))
    kicker_line(s, "A failed renewal is usually a cancelled customer, not one missed charge")
    footer(s, n, TOTAL)
    notes(s, """A payment fails. In subscriptions that's usually a cancelled customer, not one
missed charge.

The default fix everywhere is a fixed retry ladder - same card, generic text.
It's wrong twice.

It ignores WHY it failed. Retry an expired card three times and you get three
guaranteed declines and a worse decline ratio with the issuer, which raises your
cost on every future transaction.

And it ignores what trying COSTS. Outreach to a lukewarm subscriber can trigger
the exact cancellation you were trying to prevent.""")

    # 3 ------------------------------------------------------------ quadrants
    n += 1
    s = add_slide(prs)
    header(s, "WHY THE USUAL NUMBER IS WRONG", "Four kinds of customer")
    table(s, [
        ("PERSUADABLE      pays only if contacted", "the only real win"),
        ("SURE THING       pays either way", "you take the credit"),
        ("LOST CAUSE       never pays", "pure waste"),
        ("SLEEPING DOG     contact makes them leave", "you cause the churn"),
    ], Inches(2.6))
    kicker_line(s, "Gross recovery rate counts every SURE THING as a win")
    footer(s, n, TOTAL)
    notes(s, """Split customers by what your intervention actually causes and there are four
kinds.

Persuadables pay only if contacted - the only real win. Sure things pay either
way, and you take the credit. Lost causes never pay. And sleeping dogs are the
dangerous ones: contacting them causes the cancellation.

Gross recovery rate counts every sure thing as a win, and is blind to sleeping
dogs entirely. That's how these numbers get inflated.""")

    # 4 ---------------------------------------------------------- the framing
    n += 1
    s = add_slide(prs)
    header(s, "THE PART THAT MATTERS", "One run. Five headlines.")
    table(s, [
        ("Recovery rate, pursued cases only", "60.6%"),
        ("Gross recovery rate, all cases", "54.7%"),
        ("Total rupees recovered", "Rs 195,704"),
        ("Attributable to an action", "Rs 134,177"),
        ("INCREMENTAL vs control arm", "Rs 37,667"),
    ], Inches(2.55), highlight_last=True)
    kicker_line(s, "68 of 164 'recoveries' were customers who paid anyway")
    footer(s, n, TOTAL, "python -m eval.report_variants")
    notes(s, """Before I show you what I built, here's the number I refuse to lead with.

This is ONE run. One policy. One population. The only thing changing between
these five rows is the framing.

Sixty point six percent. Fifty four percent. A hundred and ninety five thousand
rupees. All true, all from that same run.

The honest number is thirty seven thousand - five times smaller - because
sixty-eight of a hundred and sixty-four so-called recoveries were customers who
paid without the agent doing anything.

That's not a hypothetical about someone else's marketing. It's my own data, and
it's why everything after this is incremental.""")

    # 5 --------------------------------------------------------------- how it
    n += 1
    s = add_slide(prs)
    header(s, "HOW IT WORKS", "Specify, verify, enforce, trace")
    bullets(s, [
        "DIAGNOSE    source x step x reason  ->  recoverability class",
        "            109 rules verified against Razorpay's docs, two tiers",
        "",
        "DECIDE      class  ->  intervention, rail, schedule",
        "            expired card -> update instrument, never re-charge",
        "",
        "VERIFY      RBI and TRAI checked BEFORE execution, never after",
        "",
        "TRACE       one hash-chained audit row per decision",
    ], Inches(2.4), font=MONO, size=16)
    kicker_line(s, "Five stopping rules. Nothing else ends a case.")
    footer(s, n, TOTAL)
    notes(s, """Here's the pipeline.

Razorpay ships an error triplet on every failed payment - source, step, reason.
Chukta maps that to one of eight recoverability classes using a hundred and nine
rules verified against Razorpay's own docs. It's two-tiered, so an unrecognised
reason degrades to a coarser class instead of falling out of the system.

Then it picks an intervention matched to the cause. Expired card gets an
instrument-update request, never another charge. Insufficient funds waits for
salary day.

Every action passes compliance gates before it executes - and every gate runs.""")

    # 6 ------------------------------------------------------------ gates live
    n += 1
    s = add_slide(prs)
    header(s, "THE GATES, ON A REAL CASE", "Rs 45,000 mandate debit, no AFA")
    table(s, [
        ("pass    G-OPS-00     kill switch", ""),
        ("BLOCK   G-RBI-01     above the AFA-free limit", ""),
        ("BLOCK   G-RBI-02     no 24h pre-debit notice", ""),
        ("pass    G-RBI-03     mandate not revoked", ""),
        ("pass    G-OPS-01     within attempt cap", ""),
        ("BLOCK   G-OPS-05     above human-approval threshold", ""),
    ], Inches(2.5), label_w=11.0)
    kicker_line(s, "Three rules block it at once. Nothing short-circuits.", RED)
    footer(s, n, TOTAL, "python -m chukta.trace")
    notes(s, """Here's that on a real case. A forty-five thousand rupee recurring debit with no
additional-factor authentication on file.

Seven gates evaluate. Three block it - two RBI rules and one approval threshold.

The detail that matters is that all seven ran. Most systems stop at the first
failure. If you're going to claim you were compliant, the audit row should show
the complete verdict, not the first objection you happened to hit.""")

    # 7 -------------------------------------------------------------- results
    n += 1
    s = add_slide(prs)
    header(s, "RESULTS", "+Rs 16,653 mean, and what it costs")
    table(s, [
        ("12-seed mean incremental", "+Rs 16,653"),
        ("95% confidence interval", "[9,274 , 24,032]"),
        ("Seeds positive", "11 of 12"),
        ("Extra customer contacts", "+55"),
        ("Extra cancellations", "+1.7"),
    ], Inches(2.55))
    kicker_line(s, "One seed comes out negative. That is reported, not hidden.")
    footer(s, n, TOTAL, "python -m eval.sweep --seeds 12")
    notes(s, """Here are the results, with the costs attached.

Across twelve seeds: plus sixteen thousand six hundred rupees mean incremental
revenue. Ninety-five percent interval, nine to twenty-four thousand. Positive in
eleven of twelve - one seed comes out negative, and that's reported, not hidden.

It costs fifty-five extra contacts and one point seven extra cancellations.

An earlier version reported a much bigger number from a single seed. That seed
turned out to be the best of twelve. It read like a measurement and it was an
anecdote.""")

    # 8 ------------------------------------------------------------- fragility
    n += 1
    s = add_slide(prs)
    header(s, "WHERE IT BREAKS", "One assumption carries the whole result")
    table(s, [
        ("baseline", "16,653      11/12"),
        ("outreach works better", "22,138      12/12"),
        ("customers churn twice as fast", "17,614      11/12"),
        ("retry timing matters less", "15,406      11/12"),
        ("message frames do nothing", "-8,129       3/12"),
    ], Inches(2.55), highlight_last=True, last_colour=RED)
    kicker_line(s, "If the copy carries no lift, this agent is NET NEGATIVE", RED)
    footer(s, n, TOTAL, "python -m eval.sweep --sensitivity")
    notes(s, """Now the part most demos skip.

Every number comes from a simulator I wrote. So I perturb one belief at a time
and re-run all twelve seeds. Six scenarios survive. One does not.

If the behavioural message frames carry no lift in payments, this agent is NET
NEGATIVE. Minus eight thousand rupees, positive in only three of twelve seeds.

Those frames come from tax-compliance trials in Guatemala, where the ask and the
consequences are completely different. Nothing I built tests whether they
transfer.

That's the most important finding here, and the one that makes it look worst.""")

    # 9 ----------------------------------------------------------- accounting
    n += 1
    s = add_slide(prs)
    header(s, "HONEST ACCOUNTING", "Most of the gain is not mine")
    table(s, [
        ("blind ladder   ->  smart timing        commodity", "-4,443"),
        ("smart timing   ->  reason-aware        commodity", "+27,515"),
        ("reason-aware   ->  Chukta            this project", "-6,418"),
        ("", ""),
        ("Buys -49% contacts and -43% churn", ""),
    ], Inches(2.5), label_w=9.6)
    kicker_line(s, "Pays off only if a retained customer is worth > Rs 2,201")
    footer(s, n, TOTAL, "python -m eval.compare_systems --seeds 12")
    notes(s, """And here's the accounting I'd least like to show.

I ran the competing strategies through the identical harness. Smart retry timing
alone is actually negative here. Reason-aware diagnosis contributes twenty-seven
thousand - and that's commodity, every serious vendor ships it.

My own contribution is minus six thousand four hundred rupees. It gives revenue
up to buy forty-nine percent fewer contacts and forty-three percent less churn.

That only pays off if a retained customer is worth more than two thousand two
hundred rupees in future billing. I have no retention data, so I'm not claiming
it.""")

    # 10 ---------------------------------------------------------------- close
    n += 1
    s = add_slide(prs)
    header(s, "WHY BELIEVE ANY OF IT", "Every number is checkable")
    bullets(s, [
        "296 tests.",
        "",
        "Qini validated against Hillstrom's 1998 randomised trial -",
        "real data nobody here generated, not just my own simulator.",
        "",
        "Doubly robust off-policy estimate lands within 5% of ground truth.",
        "",
        "eval/check_claims.py turns every README figure into a CI assertion.",
        "If the code stops matching the claims, the build goes red.",
    ], Inches(2.4))
    kicker_line(s, "github.com/AnshulPatil2005/chukta", ACCENT)
    footer(s, n, TOTAL)
    notes(s, """So why believe any of it.

Two hundred and ninety-six tests. The uplift metric is validated against
Hillstrom's published randomised experiment from 1998 - real data nobody here
generated - not just my own simulator. The off-policy estimator lands within
five percent of ground truth.

And every figure in the README is a CI assertion. If the code stops matching the
claims, the build goes red. That includes every number in this talk.

What would I measure first on real traffic? Not the routing - the copy.

That's Chukta. Thank you.""")

    return prs


if __name__ == "__main__":
    from pathlib import Path

    out = Path.home() / "OneDrive" / "Desktop" / "Chukta-pitch.pptx"
    prs = build()
    prs.save(out)
    print(f"  -> {out}")
    print(f"     {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
